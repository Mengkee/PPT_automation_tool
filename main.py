
from __future__ import print_function
from pptx import Presentation
from pptx.util import Inches
import argparse
import pandas as pd
import numpy as np
from datetime import date
import matplotlib.pyplot as plt
import seaborn as sns
import json

def dropSlides(slidesToKeep, prs):
    """Return a new presentation that has the correct slide subset.

    Param:
        - slidesToKeep: index of slides to keep from csv (int list)
        - prs: presentation (pptx.presentation)

    Return:
        - presentation with new slide subset

    """

    # get slides to delete
    indexesToRemove = [x for x in range(1, len(prs.slides._sldIdLst)+1) if x not in slidesToKeep]

    # subset report
    for i, slide in enumerate(prs.slides):
        # create slide dict
        id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(prs.slides._sldIdLst)}

        # iterate thorugh indexes
        if i+1 in indexesToRemove:
            # get slide id
            slide_id = slide.slide_id

            # remove slide
            prs.part.drop_rel(id_dict[slide_id][1])
            del prs.slides._sldIdLst[id_dict[slide_id][0]]

    return prs


def delete_placeholders(slide):
    for placeholder in slide.shapes.placeholders:
        if placeholder.has_text_frame and placeholder.text_frame.text == "":
            print("found one %s" % placeholder)
            sp = placeholder._sp
            sp.getparent().remove(sp)

def create_ppt(input, output, ppt_data): # report_data, chart
    """ Take the input powerpoint file and use it as the template for the output
    file.
    """

    print("一共{}张slide正在生成，其中默认第一页是Title第二页是目录".format(len(ppt_data["ppt_data"])))
    prs = Presentation(input)
    dropSlides(list(),prs)
    # Use the output from analyze_ppt to understand which layouts and placeholders
    # to use
    # Create a title slide first
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = ppt_data["ppt_data"][0]["title"]
    subtitle.text = ppt_data["ppt_data"][0]["subtitle"] + "Generated on {:%m-%d-%Y}".format(date.today())
    delete_placeholders(slide)

    catelogue_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(catelogue_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    content = slide.placeholders[2]
    lstr = ""
    title.text = "Table of Contents"

    number = [3,4,5,6,7,8]
    contents = [2,9,14,16,18,20]

    for i in np.arange(min(6, len(ppt_data["ppt_data"][1]["text"]))):
        item = slide.placeholders[number[i]]
        item.text = '0' + str(i + 1)
        item = slide.placeholders[contents[i]]
        item.text = ppt_data["ppt_data"][1]["text"][i]
    delete_placeholders(slide)


    for i in np.arange(len(ppt_data["ppt_data"]) - 2):
        content_slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[10]
        content = slide.placeholders[1]
        title.text = ppt_data["ppt_data"][i+2]["title"]
        subtitle.text = ppt_data["ppt_data"][i+2]["subtitle"]
        content.text = ppt_data["ppt_data"][i+2]["text"]
        delete_placeholders(slide)



    prs.save(output)

if __name__ == '__main__':
    # Opening JSON file
    f = open('config.json',"rb")

    # returns JSON object as
    # a dictionary
    ppt_data = json.load(f)

    # Closing file
    f.close()

    create_ppt("simple-template-markup.ppt", "output.ppt", ppt_data)
